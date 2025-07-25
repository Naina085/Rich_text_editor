#ppt uploading

from flask import Flask, request, jsonify
from flask_cors import CORS  # <-- Add this line
from pptx import Presentation
import os
import base64
from werkzeug.utils import secure_filename

from io import BytesIO
from PIL import Image

app = Flask(__name__)
CORS(app)
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def extract_pptx_content(file_path):
    prs = Presentation(file_path)
    content = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                content += f"<p>{shape.text}</p>"
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext
                base64_img = base64.b64encode(image_bytes).decode('utf-8')
                content += f'<img src="data:image/{image_format};base64,{base64_img}" style="max-width:100%;"><br>'
    
    return content

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'ppt' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['ppt']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)

    content = extract_pptx_content(file_path)
    return jsonify({'content': content})

if __name__ == '__main__':
    app.run(debug=True)
